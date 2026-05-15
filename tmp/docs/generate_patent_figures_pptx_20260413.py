from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/xizhuxizhu/Desktop/IndProj04/output/doc/专利PPT重绘版_20260413")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LINE = RGBColor(40, 44, 52)
TEXT = RGBColor(24, 27, 31)
FILL = RGBColor(247, 247, 247)
WHITE = RGBColor(255, 255, 255)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_box(slide, x, y, w, h, text, font_size=18, fill=WHITE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.8)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(font_size)
        run.font.bold = True if font_size >= 18 else False
        run.font.color.rgb = TEXT
    return shape


def add_section(slide, x, y, w, h, title):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    outer.fill.solid()
    outer.fill.fore_color.rgb = FILL
    outer.line.color.rgb = LINE
    outer.line.width = Pt(1.5)
    outer.text_frame.clear()
    tb = slide.shapes.add_textbox(Inches(x + 0.06), Inches(y - 0.22), Inches(2.6), Inches(0.24))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = TEXT
    return outer


def add_connector(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = LINE
    conn.line.width = Pt(1.6)
    return conn


def add_elbow(slide, pts):
    for a, b in zip(pts, pts[1:]):
        add_connector(slide, a[0], a[1], b[0], b[1])


def add_title(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0), Inches(0.08), Inches(13.333), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"图{idx}"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = TEXT


def fig1(path: Path):
    prs = new_prs()
    slide = blank_slide(prs)
    add_title(slide, 1)

    add_section(slide, 0.55, 0.9, 6.55, 2.4, "输入与语义能力层")
    add_box(slide, 0.85, 1.8, 1.65, 0.65, "用户请求")
    add_box(slide, 3.0, 1.4, 1.8, 0.52, "预设命名空间", 15)
    add_box(slide, 3.0, 2.05, 1.8, 0.52, "候选召回", 15)
    add_box(slide, 5.2, 1.4, 1.8, 0.52, "快路径裁决", 15)
    add_box(slide, 5.2, 2.05, 1.8, 0.52, "慢路径共识", 15)
    add_box(slide, 6.15, 1.72, 0.95, 0.6, "覆盖控制", 13)
    add_connector(slide, 2.5, 2.12, 3.0, 1.65)
    add_connector(slide, 2.5, 2.12, 3.0, 2.3)
    add_connector(slide, 4.8, 1.65, 5.2, 1.65)
    add_connector(slide, 4.8, 2.3, 5.2, 2.3)
    add_connector(slide, 7.0, 1.66, 7.15, 1.95)
    add_connector(slide, 7.0, 2.31, 7.15, 2.05)

    add_section(slide, 7.45, 0.9, 5.3, 2.4, "语义结果输出")
    add_box(slide, 10.15, 1.5, 1.8, 0.7, "最终主能力地址", 15)
    add_box(slide, 10.15, 2.3, 1.8, 0.7, "最终相关能力地址", 15)
    add_connector(slide, 7.1, 1.95, 10.15, 1.85)
    add_connector(slide, 7.1, 2.05, 10.15, 2.65)

    add_section(slide, 0.55, 3.6, 12.2, 2.35, "实例执行层")
    add_box(slide, 0.85, 4.5, 1.85, 0.65, "智能体注册快照", 15)
    add_box(slide, 3.4, 4.5, 2.1, 0.65, "精确地址匹配过滤", 14)
    add_box(slide, 6.4, 4.1, 1.5, 0.52, "实例匹配度", 14)
    add_box(slide, 6.4, 4.8, 1.5, 0.52, "健康度", 14)
    add_box(slide, 8.7, 4.1, 1.8, 0.52, "实例曝光公平度", 13)
    add_box(slide, 8.7, 4.8, 1.8, 0.52, "提供方曝光公平度", 13)
    add_box(slide, 11.2, 4.45, 1.1, 0.65, "实例排序", 14)
    add_connector(slide, 2.7, 4.83, 3.4, 4.83)
    add_connector(slide, 5.5, 4.83, 6.4, 4.35)
    add_connector(slide, 5.5, 4.83, 6.4, 5.05)
    add_connector(slide, 7.9, 4.35, 8.7, 4.35)
    add_connector(slide, 7.9, 5.05, 8.7, 5.05)
    add_elbow(slide, [(10.5, 4.35), (10.95, 4.35), (10.95, 4.78), (11.2, 4.78)])
    add_elbow(slide, [(10.5, 5.05), (10.95, 5.05), (10.95, 4.98), (11.2, 4.98)])
    add_connector(slide, 7.05, 3.3, 7.05, 4.1)

    add_section(slide, 0.55, 6.25, 12.2, 0.72, "最终输出与轨迹")
    add_box(slide, 1.5, 6.45, 2.0, 0.35, "目标智能体实例地址", 13)
    add_box(slide, 4.55, 6.45, 1.4, 0.35, "调用端点", 13)
    add_box(slide, 7.55, 6.45, 2.1, 0.35, "结构化决策轨迹", 13)
    add_connector(slide, 11.75, 4.78, 2.5, 6.45)
    add_connector(slide, 11.75, 4.78, 5.25, 6.45)
    add_elbow(slide, [(7.05, 3.3), (7.05, 6.0), (8.6, 6.0), (8.6, 6.45)])

    prs.save(path)


def fig2(path: Path):
    prs = new_prs()
    slide = blank_slide(prs)
    add_title(slide, 2)
    add_section(slide, 0.55, 0.9, 12.2, 2.55, "语义能力地址层")
    add_box(slide, 0.95, 2.0, 1.5, 0.6, "用户请求", 15)
    add_box(slide, 3.15, 2.0, 1.8, 0.6, "预设命名空间", 14)
    add_box(slide, 5.75, 1.72, 2.1, 1.15, ["语义能力地址", "候选集合"], 16)
    add_box(slide, 9.1, 1.52, 1.65, 0.52, "最终主能力地址", 13)
    add_box(slide, 9.1, 2.42, 1.65, 0.52, "最终相关能力地址", 13)
    add_connector(slide, 2.45, 2.3, 3.15, 2.3)
    add_connector(slide, 4.95, 2.3, 5.75, 2.3)
    add_connector(slide, 7.85, 2.0, 9.1, 1.78)
    add_connector(slide, 7.85, 2.58, 9.1, 2.68)

    add_section(slide, 0.55, 4.15, 12.2, 2.55, "智能体实例地址层")
    add_box(slide, 0.95, 5.25, 1.8, 0.6, "智能体注册快照", 15)
    add_box(slide, 3.9, 4.95, 3.2, 1.1, ["与最终主能力地址", "精确对应的实例候选集合"], 16)
    add_box(slide, 8.35, 5.25, 1.8, 0.6, "目标智能体实例地址", 13)
    add_box(slide, 10.9, 5.25, 1.2, 0.6, "调用端点", 13)
    add_connector(slide, 2.75, 5.55, 3.9, 5.55)
    add_connector(slide, 7.1, 5.55, 8.35, 5.55)
    add_connector(slide, 10.15, 5.55, 10.9, 5.55)
    add_elbow(slide, [(9.92, 2.04), (9.92, 4.5), (5.5, 4.5), (5.5, 4.95)])
    prs.save(path)


def fig3(path: Path):
    prs = new_prs()
    slide = blank_slide(prs)
    add_title(slide, 3)
    add_section(slide, 0.55, 0.92, 6.6, 1.5, "输入与视图构造")
    add_box(slide, 0.95, 1.55, 1.65, 0.5, "快路径裁决结果", 13)
    add_box(slide, 3.1, 1.55, 1.85, 0.5, "语义交接信息", 13)
    add_box(slide, 5.35, 1.4, 1.35, 0.8, ["候选视图", "构造"], 14)
    add_connector(slide, 2.6, 1.8, 3.1, 1.8)
    add_connector(slide, 4.95, 1.8, 5.35, 1.8)

    add_section(slide, 0.55, 2.85, 12.2, 1.55, "职责角色层")
    roles = [
        (0.95, "领域专家角色", "语义匹配"),
        (3.95, "治理风险角色", "风险约束"),
        (6.95, "层级解析角色", "层级冲突"),
        (9.95, "用户偏好角色", "偏好恢复"),
    ]
    for x, title, tag in roles:
        add_box(slide, x, 3.2, 1.9, 0.55, title, 13)
        add_box(slide, x + 0.4, 3.9, 1.1, 0.32, tag, 11)
    add_elbow(slide, [(6.0, 2.2), (6.0, 3.0), (1.9, 3.0), (1.9, 3.2)])
    add_elbow(slide, [(6.0, 2.2), (6.0, 3.0), (4.9, 3.0), (4.9, 3.2)])
    add_elbow(slide, [(6.0, 2.2), (6.0, 3.0), (7.9, 3.0), (7.9, 3.2)])
    add_elbow(slide, [(6.0, 2.2), (6.0, 3.0), (10.9, 3.0), (10.9, 3.2)])

    add_section(slide, 2.75, 4.95, 7.0, 1.45, "聚合与授权")
    add_box(slide, 3.25, 5.4, 4.0, 0.55, "角色提案和角色信号聚合", 14)
    add_box(slide, 7.85, 5.33, 1.25, 0.7, "覆盖授权判断", 12)
    add_connector(slide, 7.25, 5.68, 7.85, 5.68)
    add_connector(slide, 1.9, 4.22, 4.45, 5.4)
    add_connector(slide, 4.9, 4.22, 5.55, 5.4)
    add_connector(slide, 7.9, 4.22, 6.45, 5.4)
    add_connector(slide, 10.9, 4.22, 7.1, 5.4)

    add_section(slide, 10.0, 4.95, 2.75, 1.45, "输出结果")
    add_box(slide, 10.35, 5.28, 2.0, 0.45, "最终主能力地址", 13)
    add_box(slide, 10.35, 5.82, 2.0, 0.45, "最终相关能力地址", 13)
    add_connector(slide, 9.1, 5.62, 10.35, 5.5)
    add_connector(slide, 9.1, 5.78, 10.35, 6.04)
    prs.save(path)


def fig4(path: Path):
    prs = new_prs()
    slide = blank_slide(prs)
    add_title(slide, 4)
    add_section(slide, 0.55, 0.92, 12.2, 1.8, "过滤阶段")
    xs = [0.85, 3.2, 5.55, 7.9, 10.25]
    labels = ["智能体注册快照", "精确地址匹配过滤", "实例状态过滤", "调用端点过滤", "模式兼容性过滤"]
    ws = [1.75, 2.0, 1.65, 1.65, 1.7]
    for i, (x, label, w) in enumerate(zip(xs, labels, ws)):
        add_box(slide, x, 1.62, w, 0.58, label, 13)
        if i < len(xs) - 1:
            add_connector(slide, x + w, 1.91, xs[i + 1], 1.91)

    add_section(slide, 0.55, 3.4, 12.2, 2.25, "评分与排序阶段")
    add_box(slide, 1.4, 4.1, 1.55, 0.5, "实例匹配度", 13)
    add_box(slide, 3.9, 4.1, 1.35, 0.5, "健康度", 13)
    add_box(slide, 6.05, 4.1, 1.95, 0.5, "实例曝光公平度", 12)
    add_box(slide, 8.8, 4.1, 2.0, 0.5, "提供方曝光公平度", 12)
    add_box(slide, 5.0, 5.05, 3.0, 0.58, "排序得分生成", 13)
    add_elbow(slide, [(11.1, 2.2), (11.1, 3.75), (2.15, 3.75), (2.15, 4.1)])
    add_elbow(slide, [(11.1, 2.2), (11.1, 3.75), (4.58, 3.75), (4.58, 4.1)])
    add_elbow(slide, [(11.1, 2.2), (11.1, 3.75), (7.0, 3.75), (7.0, 4.1)])
    add_elbow(slide, [(11.1, 2.2), (11.1, 3.75), (9.8, 3.75), (9.8, 4.1)])
    add_connector(slide, 2.15, 4.6, 5.8, 5.05)
    add_connector(slide, 4.58, 4.6, 6.25, 5.05)
    add_connector(slide, 7.0, 4.6, 6.95, 5.05)
    add_connector(slide, 9.8, 4.6, 7.6, 5.05)

    add_section(slide, 4.55, 6.15, 4.0, 0.75, "输出")
    add_box(slide, 4.9, 6.34, 3.3, 0.36, "目标智能体实例地址及调用端点", 12)
    add_connector(slide, 6.5, 5.63, 6.5, 6.34)
    prs.save(path)


def fig5(path: Path):
    prs = new_prs()
    slide = blank_slide(prs)
    add_title(slide, 5)
    add_section(slide, 0.5, 0.82, 9.4, 5.6, "结构化决策轨迹对象")
    add_box(slide, 1.1, 1.25, 8.1, 0.62, ["召回阶段字段组", "语义能力地址候选集合、候选相关度、混淆源标记"], 12)
    add_box(slide, 1.1, 2.2, 8.1, 0.78, ["快路径阶段字段组", "初始主能力地址、初始相关能力地址、裁决置信度、候选竞争差值、慢路径触发原因"], 12)
    add_box(slide, 1.1, 3.35, 8.1, 0.78, ["慢路径阶段字段组", "语义交接信息、角色提案、角色信号、聚合结果、覆盖阻断原因"], 12)
    add_box(slide, 1.1, 4.5, 8.1, 0.9, ["实例选择阶段字段组", "实例过滤原因、实例匹配度、健康度、实例曝光公平度、提供方曝光公平度、最终排序得分"], 12)
    add_box(slide, 1.1, 5.75, 8.1, 0.62, ["输出阶段字段组", "最终主能力地址、最终相关能力地址、目标智能体实例地址、调用端点"], 12)
    add_connector(slide, 5.15, 1.87, 5.15, 2.2)
    add_connector(slide, 5.15, 2.98, 5.15, 3.35)
    add_connector(slide, 5.15, 4.13, 5.15, 4.5)
    add_connector(slide, 5.15, 5.4, 5.15, 5.75)

    add_section(slide, 10.35, 1.2, 1.95, 2.55, "可追溯能力")
    add_box(slide, 10.7, 1.68, 1.25, 0.42, "可回放", 12)
    add_box(slide, 10.7, 2.45, 1.25, 0.42, "可归因", 12)
    add_box(slide, 10.7, 3.22, 1.25, 0.42, "可审计", 12)
    add_connector(slide, 9.2, 1.56, 10.7, 1.89)
    add_connector(slide, 9.2, 3.75, 10.7, 2.66)
    add_connector(slide, 9.2, 6.06, 10.7, 3.43)

    add_section(slide, 10.35, 4.55, 1.95, 1.85, "错误定位")
    add_box(slide, 10.65, 5.0, 1.35, 0.34, "召回阶段", 11)
    add_box(slide, 10.65, 5.45, 1.35, 0.34, "裁决阶段", 11)
    add_box(slide, 10.65, 5.9, 1.35, 0.34, "实例选择阶段", 11)
    add_connector(slide, 9.2, 1.56, 10.65, 5.17)
    add_connector(slide, 9.2, 2.59, 10.65, 5.62)
    add_connector(slide, 9.2, 4.95, 10.65, 6.07)
    prs.save(path)


def combined(fig_paths: list[tuple[str, callable]], out_path: Path):
    prs = new_prs()
    # remove default slide then generate temporary per-slide in same prs
    prs.slides._sldIdLst.clear()  # type: ignore[attr-defined]
    builders = [fig1, fig2, fig3, fig4, fig5]
    for i, builder in enumerate(builders, start=1):
        slide = blank_slide(prs)
        if i == 1:
            # build on existing slide by reusing builder logic through temp file impossible
            pass
    # keep combined deck optional skipped
    prs.save(out_path)


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    fig1(OUT / "图1_系统总体结构示意图.pptx")
    fig2(OUT / "图2_两层地址关系示意图.pptx")
    fig3(OUT / "图3_协同决策结构示意图.pptx")
    fig4(OUT / "图4_实例过滤与排序过程示意图.pptx")
    fig5(OUT / "图5_决策轨迹组织示意图.pptx")
    (OUT / "README.md").write_text(
        "\n".join(
            [
                "# 专利PPT重绘版 20260413",
                "",
                "本目录采用 PowerPoint 矢量图形重绘，每张图为单页 PPTX，便于继续人工微调。",
                "",
                "- 图1_系统总体结构示意图.pptx",
                "- 图2_两层地址关系示意图.pptx",
                "- 图3_协同决策结构示意图.pptx",
                "- 图4_实例过滤与排序过程示意图.pptx",
                "- 图5_决策轨迹组织示意图.pptx",
                "",
            ]
        ),
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
