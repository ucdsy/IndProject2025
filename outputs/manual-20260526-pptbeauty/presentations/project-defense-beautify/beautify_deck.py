from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


INPUT = Path(
    "/Users/xizhuxizhu/Downloads/%E9%87%8D%E7%82%B9%E9%A1%B9%E7%9B%AE%E9%AA%8C%E6%94%B6%E7%AD%94%E8%BE%A9PPT_%E6%9D%90%E6%96%99%E5%9B%BE%E7%89%88_v1_20260526.pptx"
)
OUTPUT = Path(
    "/Users/xizhuxizhu/Downloads/重点项目验收答辩PPT_材料图版_美化版_20260526.pptx"
)

EMU = 914400


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def inches(value: int | float) -> float:
    return value / EMU


def set_text(shape, text: str, *, size: float | None = None, color: str = "102033", bold: bool | None = None):
    tf = shape.text_frame
    old_size = None
    old_bold = None
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                old_size = r.font.size
            if r.font.bold is not None:
                old_bold = r.font.bold
            break
        if old_size is not None or old_bold is not None:
            break
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size) if size is not None else (old_size or Pt(10))
    r.font.bold = old_bold if bold is None else bold
    r.font.color.rgb = rgb(color)


def soften_shape(shape, fill="f7fbff", line="d9e8f7"):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.9)


def remove_black_material_bases(slide):
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        x, y, w, h = map(inches, (shape.left, shape.top, shape.width, shape.height))
        if x < 0.45 and 0.95 < y < 1.2 and w > 12.0 and h > 5.8:
            soften_shape(shape, fill="f7fbff", line=None)
        elif x < 0.75 and 1.2 < y < 1.6 and w > 7.0 and h > 3.8:
            soften_shape(shape, fill="ffffff", line="d9e8f7")


def polish_text(slide_no: int, slide):
    replacements = {
        "PPT 负责把“已完成什么、为什么有技术含量、凭什么可验收”讲成一条证据链。":
            "本页说明项目如何从任务书要求落到可验收证据链。",
        "材料原图用于主讲页：完整展示候选集合构造、结构化判别、协作复核、共识授权、执行映射和过程记录。":
            "技术路线贯通候选集合构造、结构化判别、协作复核、授权控制、执行映射和过程记录。",
        "材料图源：图 2 两层地址关系":
            "证据图说明：两层地址关系",
        "PPT处理建议":
            "答辩表达",
        "原图是黑白流程图，不建议全屏主视觉；本版把它作为材料证据图，右侧用三点归纳替代长流程。":
            "原图作为证据图，右侧用三点归纳说明能力地址与实例地址的分层关系。",
        "采用材料中的两层地址关系图，说明“能力地址”和“执行实例”分层。":
            "能力地址与实例地址分层，支撑路由归因和执行约束。",
        "材料原图用于主讲页：四类角色并行复核、提案聚合、第二轮定向复核和授权检查构成受控改判闭环。":
            "四类角色并行复核、提案聚合、二轮定向复核和授权检查共同构成受控改判闭环。",
        "材料图源：结构化决策轨迹":
            "证据图说明：结构化决策轨迹",
        "采用材料中的“结构化决策轨迹”原图，右侧压缩为可信留痕验收表述。":
            "结构化决策轨迹支撑候选、裁决、复核、授权和执行映射全程留痕。",
        "系统演示建议：使用真实截图替换此处模拟界面，答辩时重点讲“候选—裁决—复核—执行—留痕”。":
            "现场演示重点：候选—裁决—复核—执行—留痕。",
        "材料原图：原型系统运行截图":
            "原型系统运行截图",
        "本页替换为成果汇编中的真实原型截图，比模拟界面更适合作为验收证据。":
            "真实原型截图支撑现场演示，展示候选、裁决、复核、执行和留痕闭环。",
        "加入材料图 10，作为样本扩展后排序稳定的诊断证据。":
            "累计样本池稳定性用于说明方法排序没有被新增复杂样本推翻。",
        "采用成果汇编主结果图；右侧只保留答辩时要讲的四个结论。":
            "结构化判别提供主要增益，协作复核在复杂边界样本上形成受控修正。",
        "采用成果汇编配对消融图，支撑“协作收益来自职责差异 + 授权控制”。":
            "配对消融证明协作收益来自职责差异与授权控制。",
        "案例页建议现场只讲一个样本，重点证明“可回放、可归因、可解释”。":
            "围绕一个样本展示误判原因、复核触发、授权改判和过程回放。",
        "证据链：成果汇编管总，技术报告讲机制，论文提供实验硬证据，PPT 汇报“任务书—成果—验收”的对应关系。":
            "证据链：成果汇编统筹材料，技术报告说明机制，论文提供实验硬证据，原型系统支撑现场验证。",
    }
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        old = " ".join(shape.text.strip().split())
        if old in replacements:
            size = 10.5
            if slide_no in {5, 9, 11, 13, 14} and "材料原图" in old:
                size = 10.0
            set_text(shape, replacements[old], size=size)


def polish_cards(slide):
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        x, y, w, h = map(inches, (shape.left, shape.top, shape.width, shape.height))
        if w > 3.5 and h > 0.5 and y > 5.6:
            try:
                soften_shape(shape, fill="eaf3ff", line="8fbced")
            except Exception:
                pass


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def remove_material_overlay(slide):
    shapes = list(slide.shapes)
    start = None
    for idx, shape in enumerate(shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        x, y, w, h = map(inches, (shape.left, shape.top, shape.width, shape.height))
        if x < 0.45 and 0.95 < y < 1.2 and w > 12.0 and h > 5.8:
            start = idx
            break
    if start is None:
        return
    for shape in reversed(shapes[start:]):
        remove_shape(shape)


def main():
    prs = Presentation(INPUT)
    for i, slide in enumerate(prs.slides, 1):
        remove_black_material_bases(slide)
        polish_cards(slide)
        polish_text(i, slide)
        if i in {5, 7, 9, 10, 11, 12, 13, 14}:
            remove_material_overlay(slide)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
