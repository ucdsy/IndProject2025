from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/xizhuxizhu/Desktop/IndProj04")
INPUT = Path("/Users/xizhuxizhu/Downloads/重点项目验收答辩PPT_生成图模板融合版_20260526.pptx")
OUTPUT = Path("/Users/xizhuxizhu/Downloads/重点项目验收答辩PPT_论文图融合美化版_20260526.pptx")

FIG1 = ROOT / "output/doc/gjtx_submission_20260413/figures/fig1_framework.png"
FIG2 = ROOT / "output/doc/gjtx_submission_20260413/figures/fig2_review_flow.png"
DATA_FIG = ROOT / "output/doc/gjtx_submission_20260413/figures/09_historical_cumulative_train_test.png"

EMU = 914400
SLIDE_W, SLIDE_H = 13.333333, 7.5

C = {
    "navy": "061a46",
    "blue": "0b63ce",
    "blue2": "1e83e6",
    "cyan": "26b6ff",
    "pale": "f7fbff",
    "pale2": "eaf3ff",
    "line": "c8dbef",
    "ink": "102033",
    "muted": "53657a",
    "orange": "f59e0b",
    "green": "16a35c",
    "purple": "6b49c8",
    "red": "b93232",
    "white": "ffffff",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def emu_to_in(v: int | float) -> float:
    return v / EMU


def pos(x, y, w, h):
    return Inches(x), Inches(y), Inches(w), Inches(h)


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def send_to_back(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)
    parent.insert(2, shape._element)


def clear_slide(slide):
    for shape in reversed(list(slide.shapes)):
        remove_shape(shape)


def add_box(slide, x, y, w, h, fill="ffffff", line="c8dbef", radius=True, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        *pos(x, y, w, h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C.get(fill, fill))
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(C.get(line, line))
        shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w, h, size=12, color="ink", bold=False, align="left", fill=None, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        *pos(x, y, w, h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(C.get(fill, fill))
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(C.get(line, line))
        shape.line.width = Pt(0.9)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(C.get(color, color))
    return shape


def add_fit_picture(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw = iw * scale
    ph = ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), *pos(px, py, pw, ph))


def draw_base(slide, section: str, title: str, subtitle: str, page: int):
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill="f7fbff", line=None, radius=False)
    add_box(slide, 0, 0, SLIDE_W, 0.16, fill="0b63ce", line=None, radius=False)
    add_box(slide, 0, 0.16, SLIDE_W, 0.06, fill="26b6ff", line=None, radius=False)
    add_box(slide, 0.4, 0.34, 0.46, 0.46, fill="0b63ce", line="0b63ce", radius=True)
    add_text(slide, section, 0.4, 0.44, 0.46, 0.15, size=12, color="white", bold=True, align="center")
    add_text(slide, title, 0.98, 0.33, 9.7, 0.32, size=20, color="navy", bold=True)
    add_text(slide, subtitle, 1.0, 0.75, 9.9, 0.16, size=8.2, color="muted")
    add_text(slide, "CNNIC", 11.56, 0.27, 1.1, 0.26, size=20, color="blue", bold=True, align="right")
    add_text(slide, "中国互联网络信息中心", 10.9, 0.56, 1.75, 0.14, size=7.5, color="muted", align="right")
    add_text(slide, "面向互联网基础资源的大模型多智能体协作与可信认知标识技术研究", 0.52, 7.14, 5.6, 0.14, size=6.2, color="muted")
    add_text(slide, f"{page:02d}", 12.35, 7.12, 0.26, 0.13, size=7, color="muted", align="right")


def set_text_if_contains(slide, contains: str, replacement: str, size: float | None = None):
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text:
            continue
        if contains not in shape.text:
            continue
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = replacement
        r.font.name = "Microsoft YaHei"
        r.font.bold = True if len(replacement) < 70 else False
        if size:
            r.font.size = Pt(size)
        r.font.color.rgb = rgb(C["ink"])


def normalize_fonts(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    if run.font.size and run.font.size.pt < 6.5:
                        run.font.size = Pt(6.5)


def remove_material_overlay(slide):
    shapes = list(slide.shapes)
    start = None
    for idx, shape in enumerate(shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        x, y, w, h = map(emu_to_in, (shape.left, shape.top, shape.width, shape.height))
        if x < 0.45 and 0.95 < y < 1.2 and w > 12.0 and h > 5.8:
            start = idx
            break
    if start is None:
        return
    for shape in reversed(shapes[start:]):
        remove_shape(shape)


def cover_content_area(slide):
    bg = add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill="f7fbff", line=None, radius=False)
    send_to_back(bg)
    add_box(slide, 0.35, 0.96, 12.63, 6.12, fill="f7fbff", line=None, radius=False)
    add_box(slide, 0.54, 1.14, 12.24, 5.52, fill="ffffff", line="c8dbef", radius=True)


def paper_framework_slide(slide):
    clear_slide(slide)
    draw_base(
        slide,
        "二",
        "2.1 总体架构：论文架构图支撑完整技术路线",
        "从请求输入到候选构造、结构化判别、协作复核、授权控制、执行映射和过程记录，全链路可验证。",
        5,
    )
    add_box(slide, 0.48, 1.02, 12.38, 5.82, fill="ffffff", line="c8dbef", radius=True)
    add_text(slide, "论文图 1｜总体技术框架", 0.84, 1.28, 2.2, 0.24, size=11, color="blue", bold=True)
    add_fit_picture(slide, FIG1, 0.76, 1.58, 11.78, 4.34)
    add_box(slide, 0.8, 6.08, 11.72, 0.42, fill="eaf3ff", line="8fbced", radius=True)
    add_text(
        slide,
        "答辩讲法：能力命名约束候选边界，结构化语义路由完成候选内裁决，职责化复核处理复杂冲突，可信轨迹支撑回放核验。",
        1.04,
        6.18,
        10.9,
        0.16,
        size=9,
        color="ink",
        bold=True,
    )


def paper_review_slide(slide):
    clear_slide(slide)
    draw_base(
        slide,
        "二",
        "2.5 关键技术三：论文复核架构支撑职责化协作与授权改判",
        "协作触发、角色复核、提案聚合、二轮定向复核和授权检查分离，避免无约束多轮讨论。",
        9,
    )
    add_box(slide, 0.48, 1.02, 12.38, 5.82, fill="ffffff", line="c8dbef", radius=True)
    add_text(slide, "论文图 2｜职责化协作复核流程", 0.84, 1.28, 2.45, 0.24, size=11, color="orange", bold=True)
    add_fit_picture(slide, FIG2, 0.74, 1.54, 11.86, 4.26)
    add_box(slide, 0.8, 6.08, 11.72, 0.42, fill="fff7ed", line="f3c37c", radius=True)
    add_text(
        slide,
        "答辩讲法：四类角色不是增加模型调用次数，而是把任务匹配、风险治理、层级冲突和用户偏好拆开形成异质证据。",
        1.04,
        6.18,
        10.9,
        0.16,
        size=9,
        color="ink",
        bold=True,
    )


def data_volume_slide(slide):
    clear_slide(slide)
    draw_base(
        slide,
        "三",
        "3.2 数据规模与稳定性：按累计样本量呈现方法效果",
        "样本池从小规模逐步扩大到 train=450、test=113，方法排序保持稳定，支撑结论不是偶然样例。",
        12,
    )
    add_box(slide, 0.48, 1.02, 12.38, 5.82, fill="ffffff", line="c8dbef", radius=True)
    metrics = [
        ("563", "冻结样本", "blue"),
        ("450", "train", "green"),
        ("113", "test", "orange"),
        ("0.929", "测试侧扩展配置终值", "red"),
    ]
    for i, (v, label, col) in enumerate(metrics):
        x = 0.78 + i * 1.62
        add_box(slide, x, 1.28, 1.35, 0.68, fill="ffffff", line=C[col], radius=True)
        add_text(slide, v, x + 0.08, 1.4, 1.16, 0.2, size=16, color=col, bold=True, align="center")
        add_text(slide, label, x + 0.08, 1.68, 1.16, 0.12, size=6.8, color="muted", align="center")
    add_box(slide, 0.74, 2.16, 8.2, 3.72, fill="ffffff", line="d9e8f7", radius=True)
    add_fit_picture(slide, DATA_FIG, 0.92, 2.33, 7.84, 3.36)
    cards = [
        ("样本扩大", "从小样本扩展到 563 条冻结样本，训练/测试分开统计。", "blue"),
        ("排序稳定", "结构化语义路由、默认协作和扩展配置长期高于规则基线。", "green"),
        ("测试终值", "测试侧扩展配置最终约 0.929，默认协作约 0.885。", "red"),
    ]
    for i, (title, body, col) in enumerate(cards):
        y = 2.2 + i * 1.1
        add_box(slide, 9.2, y, 3.25, 0.86, fill="ffffff", line=C[col], radius=True)
        add_text(slide, f"0{i+1}", 9.36, y + 0.28, 0.38, 0.16, size=8.4, color=col, bold=True, align="center")
        add_text(slide, title, 9.86, y + 0.18, 1.0, 0.14, size=9.4, color=col, bold=True)
        add_text(slide, body, 9.86, y + 0.46, 2.18, 0.18, size=7.2, color="muted")
    add_box(slide, 0.8, 6.14, 11.72, 0.38, fill="eaf3ff", line="8fbced", radius=True)
    add_text(slide, "结论：随着样本池扩展，方法排序没有被新增复杂样本推翻，实验结论具备稳定性支撑。", 1.04, 6.23, 10.8, 0.14, size=9, color="ink", bold=True)


def cleanup_internal_hints(slide):
    replacements = {
        "PPT处理建议": "答辩表达",
        "原图是黑白流程图，不建议全屏主视觉；本版把它作为材料证据图，右侧用三点归纳替代长流程。": "作为备份证据图，用于解释能力地址与实例地址的分层边界。",
        "采用材料中的两层地址关系图，说明“能力地址”和“执行实例”分层。": "能力地址回答“哪类能力处理”，实例地址回答“哪个智能体执行”。",
        "材料图源": "证据图说明",
        "材料原图": "证据图",
        "系统演示建议": "演示重点",
        "案例页建议": "案例讲法",
        "本页替换为成果汇编中的真实原型截图，比模拟界面更适合作为验收证据。": "真实原型截图支撑现场演示，展示候选、裁决、复核、执行和留痕闭环。",
    }
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text:
            continue
        for old, new in replacements.items():
            if old in shape.text:
                shape.text = shape.text.replace(old, new)


def main():
    prs = Presentation(INPUT)
    normalize_fonts(prs)
    cleanup_internal_hints(prs.slides[6])
    remove_material_overlay(prs.slides[6])
    paper_framework_slide(prs.slides[4])
    paper_review_slide(prs.slides[8])
    data_volume_slide(prs.slides[11])
    for idx in range(len(prs.slides)):
        cleanup_internal_hints(prs.slides[idx])
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
