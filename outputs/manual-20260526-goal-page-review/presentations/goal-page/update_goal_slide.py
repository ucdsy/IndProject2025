from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPT = Path("/Users/xizhuxizhu/Desktop/重点项目验收答辩PPT.pptx")
BACKUP = Path("/Users/xizhuxizhu/Desktop/重点项目验收答辩PPT_目标页修改前备份_20260526.pptx")

EMU = 914400
C = {
    "navy": "061a46",
    "blue": "0b63ce",
    "blue2": "1e83e6",
    "cyan": "26b6ff",
    "pale": "f7fbff",
    "pale2": "eaf3ff",
    "line": "c8dbef",
    "ink": "102033",
    "muted": "53657a",
    "orange": "f59e0b",
    "green": "16a35c",
    "purple": "6b49c8",
    "red": "b93232",
    "white": "ffffff",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def pos(x, y, w, h):
    return Inches(x), Inches(y), Inches(w), Inches(h)


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def add_box(slide, x, y, w, h, fill="ffffff", line="c8dbef", radius=True, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        *pos(x, y, w, h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C.get(fill, fill))
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(C.get(line, line))
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=12,
    color="ink",
    bold=False,
    align="left",
    fill=None,
    line=None,
    radius=False,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        *pos(x, y, w, h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(C.get(fill, fill))
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(C.get(line, line))
        shape.line.width = Pt(0.9)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(C.get(color, color))
    return shape


def set_existing_text(shape, text: str, *, size=18, color="blue", bold=True, width=None, page_no=False):
    if width is not None:
        shape.width = Inches(width)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT if page_no else PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(C.get(color, color))


def add_goal_card(slide, x, y, w, h, num, title, body, color):
    add_box(slide, x, y, w, h, fill="ffffff", line=C[color], radius=True)
    add_box(slide, x + 0.18, y + 0.18, 0.34, 0.34, fill=C[color], line=C[color], radius=True)
    add_text(slide, num, x + 0.18, y + 0.25, 0.34, 0.1, size=8, color="white", bold=True, align="center")
    add_text(slide, title, x + 0.62, y + 0.18, w - 0.82, 0.18, size=10.5, color=color, bold=True)
    add_text(slide, body, x + 0.24, y + 0.62, w - 0.42, h - 0.78, size=7.8, color="muted")


def add_flow_node(slide, x, y, text, color):
    add_box(slide, x, y, 1.62, 0.42, fill="ffffff", line=C[color], radius=True)
    add_text(slide, text, x + 0.06, y + 0.13, 1.5, 0.1, size=7.4, color=color, bold=True, align="center")


def main():
    if not BACKUP.exists():
        shutil.copy2(PPT, BACKUP)

    prs = Presentation(PPT)
    slide = prs.slides[3]

    # Remove the previous non-editable target illustration.
    for shape in reversed(list(slide.shapes)):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            x, y, w, h = shape.left / EMU, shape.top / EMU, shape.width / EMU, shape.height / EMU
            if y > 1.0 and w > 5:
                remove_shape(shape)

    # Update title and page number.
    set_existing_text(
        slide.shapes[4],
        "1.2 研究目标：构建面向互联网基础资源场景的多智能体可信能力发现与路由验证框架",
        size=18,
        color="blue",
        bold=True,
        width=9.8,
    )
    set_existing_text(slide.shapes[8], "04", size=8, color="muted", bold=False, page_no=True)

    # Content panel.
    add_box(slide, 0.58, 1.05, 12.15, 5.78, fill="ffffff", line="c8dbef", radius=True, transparency=0)
    add_text(
        slide,
        "以智能体能力命名与语义路由为典型验证载体，将自然语言请求转化为可复核、可授权、可留痕的能力调用决策。",
        0.95,
        1.28,
        10.7,
        0.28,
        size=12.4,
        color="ink",
        bold=True,
        align="center",
    )

    # Goal flow.
    flow = [
        ("能力命名空间", "blue"),
        ("候选内判别", "green"),
        ("职责化复核", "orange"),
        ("授权改判", "purple"),
        ("可信留痕", "red"),
        ("执行映射", "blue2"),
    ]
    for i, (label, color) in enumerate(flow):
        x = 1.05 + i * 1.9
        add_flow_node(slide, x, 2.02, label, color)
        if i < len(flow) - 1:
            add_text(slide, "→", x + 1.62, 2.12, 0.22, 0.1, size=12, color="muted", bold=True, align="center")

    # Five concrete objectives.
    cards = [
        ("01", "能力对象可组织", "构建智能体能力命名空间，明确能力地址与实例地址边界。", "blue"),
        ("02", "请求意图可判别", "在固定候选集合内输出主能力、相关能力、置信度和竞争候选说明。", "green"),
        ("03", "复杂样本可复核", "通过任务匹配、治理风险、层级解析、用户偏好四类角色进行职责化协作。", "orange"),
        ("04", "改判过程可控制", "将协作触发与最终改判分离，通过票数、证据强度、共识增益和风险确认进行显式授权。", "purple"),
        ("05", "判断过程可追溯", "记录候选快照、结构化裁决、角色提案、授权判断、最终能力地址和执行实例映射。", "red"),
    ]
    for i, card in enumerate(cards):
        add_goal_card(slide, 0.85 + i * 2.36, 3.03, 2.05, 1.68, *card)

    add_box(slide, 0.86, 5.24, 11.55, 0.56, fill="eaf3ff", line="8fbced", radius=True)
    add_text(
        slide,
        "形成“能力命名空间—候选内判别—职责化复核—授权改判—可信留痕—执行映射”的可验证技术链路。",
        1.1,
        5.42,
        10.85,
        0.15,
        size=11.0,
        color="blue",
        bold=True,
        align="center",
    )

    add_box(slide, 0.86, 6.06, 11.55, 0.42, fill="071a46", line="071a46", radius=True)
    add_text(
        slide,
        "验收口径：目标不是泛泛做 Agent 应用，而是把多智能体协作与可信认知标识落到可评测、可回放、可交付的原型链路。",
        1.1,
        6.18,
        10.85,
        0.12,
        size=9.0,
        color="white",
        bold=True,
        align="center",
    )

    prs.save(PPT)
    print(PPT)
    print(BACKUP)


if __name__ == "__main__":
    main()
